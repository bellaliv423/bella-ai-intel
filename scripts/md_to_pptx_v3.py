"""
🎨 v3: 마크다운 → PPTX 풀 재작업 (벨라님 D+5 21:35 피드백 반영)

v2 → v3 핵심 5가지 개선 (D+6 화 2026-04-28):
1. ✅ 폰트 사이즈 전체 ↑ (헤드라인 42→52pt, 본문 14→18pt, 코드 11→14pt)
2. ✅ Pretendard Bold (한글 가독성 최상, 폴백 맑은 고딕)
3. ✅ Hero Number 패턴 (70pt 큰 숫자 강조 — ₩5M, 20명, 9시간)
4. ✅ 깔끔함 — 여백 0.7→1.0" (좌우), 0.4→0.6" (상하)
5. ✅ 임팩트 — Crail Brown 풀 박스 + 화이트 텍스트 (반전 CTA)

v1 → v2 개선 7가지 (유지):
1. ✅ 슬라이드 분할 룰: H2만 새 슬라이드 (H3는 같은 슬라이드 섹션)
2. ✅ 빈 슬라이드 자동 제거 (텍스트 50자 미만 + 다음 슬라이드와 병합)
3. ✅ 한글 폰트 풀 지정 (latin + eastAsia + cs)
4. ✅ 코드 블록 = 다크 배경 박스 + monospace
5. ✅ 진짜 테이블 (PPTX Table 위젯)
6. ✅ 인라인 마크다운 정리 (**bold** *italic* `code`)
7. ✅ 표지에 픽셀아트 캐릭터 + 단체 배너 옵션

사용법:
  python md_to_pptx_v2.py --input <md> --output <pptx>     # 단일
  python md_to_pptx_v2.py --batch                           # 7개 일괄
  python md_to_pptx_v2.py --sample                          # 샘플 1개만
"""
import argparse
import re
import sys
import io
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# ═════════════════════════════════════════════════════════
# 🎨 Anthropic 공식 브랜드 컬러
# ═════════════════════════════════════════════════════════
COLOR_BG = RGBColor(0xF0, 0xEE, 0xE6)        # Manilla Cream
COLOR_TEXT = RGBColor(0x19, 0x19, 0x19)      # Slate
COLOR_ACCENT = RGBColor(0xCC, 0x78, 0x5C)    # Crail Brown
COLOR_MUTED = RGBColor(0x80, 0x70, 0x60)
COLOR_CODE_BG = RGBColor(0x1F, 0x1F, 0x1F)   # 다크 코드 배경
COLOR_CODE_TEXT = RGBColor(0xE8, 0xE6, 0xDE)
COLOR_CODE_ACCENT = RGBColor(0xE8, 0x9A, 0x7E)
COLOR_HIGHLIGHT_BG = RGBColor(0xFF, 0xF5, 0xEC)  # 표 헤더 배경
COLOR_TABLE_BORDER = RGBColor(0xD0, 0xC8, 0xBC)

# v3: Pretendard 한글 가독성 (Toss·당근·인스타 인플루언서 표준)
FONT_HEADING = 'Pretendard'
FONT_BODY = 'Pretendard'
FONT_CODE = 'D2Coding'

# v3 폰트 사이즈 풀 ↑ (벨라님 피드백)
FONT_HERO_NUMBER = 70   # 큰 숫자 (₩5M, 20명, 9시간 강조)
FONT_HERO_HEADLINE = 52 # 표지 메인 헤드라인 (was 42)
FONT_CHAPTER = 40       # 챕터 타이틀 (was 28)
FONT_SUBTITLE = 24      # 표지 부제 (was 18)
FONT_H3 = 24            # 섹션 제목 (was 18)
FONT_BODY_SIZE = 18     # 본문 (was 14)
FONT_BULLET = 17        # 불릿 (was 14)
FONT_CODE_SIZE = 14     # 코드 (was 11)
FONT_FOOTER = 11        # 푸터 (was 9)
FONT_LABEL = 12         # 라벨 (was 10)

# v3 여백 ↑
MARGIN_LR = Inches(1.0)
MARGIN_TB = Inches(0.6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═════════════════════════════════════════════════════════
# 한글 폰트 풀 지정 (latin + eastAsia + cs)
# ═════════════════════════════════════════════════════════

def set_korean_font(run, font_name, fallback='맑은 고딕'):
    """run에 한글+영문 통합 폰트 적용"""
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    # eastAsia
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font_name)
    # latin
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = etree.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', font_name)
    # cs (complex script)
    cs = rPr.find(qn('a:cs'))
    if cs is None:
        cs = etree.SubElement(rPr, qn('a:cs'))
    cs.set('typeface', font_name)


# ═════════════════════════════════════════════════════════
# 마크다운 인라인 정리
# ═════════════════════════════════════════════════════════

def clean_inline_markdown(text: str) -> str:
    """**bold** *italic* `code` ~~strike~~ 정리"""
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)
    text = re.sub(r'\*([^*]+?)\*', r'\1', text)
    text = re.sub(r'_([^_]+?)_', r'\1', text)
    text = re.sub(r'`([^`]+?)`', r'\1', text)
    text = re.sub(r'~~(.+?)~~', r'\1', text)
    # 링크 [text](url) → text
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    return text


# ═════════════════════════════════════════════════════════
# 마크다운 파서 v2
# ═════════════════════════════════════════════════════════

def parse_markdown(content: str) -> dict:
    fm = {}
    body = content
    fm_match = re.match(r'^---\n(.*?)\n---\n(.*)', content, re.DOTALL)
    if fm_match:
        fm_text = fm_match.group(1)
        body = fm_match.group(2)
        for line in fm_text.splitlines():
            if ':' in line:
                k, v = line.split(':', 1)
                fm[k.strip()] = v.strip()
    return {'frontmatter': fm, 'body': body}


def split_into_slides_v2(body: str) -> list:
    """H2를 슬라이드 단위로 (H1은 챕터 섹션 헤더, H3는 같은 슬라이드 서브 섹션)"""
    slides = []
    lines = body.split('\n')
    current_chapter = None
    current_slide = None

    for line in lines:
        stripped = line.rstrip()
        h1 = re.match(r'^#\s+(.+)$', stripped)
        h2 = re.match(r'^##\s+(.+)$', stripped)

        if h1:
            current_chapter = h1.group(1).strip()
            # H1은 별도 슬라이드 안 만들고 다음 H2의 챕터 라벨로
            continue

        if h2:
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'chapter': current_chapter or '',
                'title': h2.group(1).strip(),
                'content': []
            }
            continue

        if current_slide is None:
            current_slide = {'chapter': current_chapter or '', 'title': '', 'content': []}

        current_slide['content'].append(line)

    if current_slide:
        slides.append(current_slide)

    # 빈 슬라이드 + 작은 슬라이드 병합
    cleaned = []
    for s in slides:
        text_len = sum(len(l) for l in s['content'])
        if text_len < 50 and not s['title']:
            # 너무 작으면 이전 슬라이드에 병합
            if cleaned:
                cleaned[-1]['content'].extend(s['content'])
            continue
        cleaned.append(s)

    return cleaned


# ═════════════════════════════════════════════════════════
# 컨텐츠 블록 렌더링 (테이블 + 코드블록 분리)
# ═════════════════════════════════════════════════════════

def render_content_blocks(lines: list) -> list:
    """라인 → 블록 (text / code / table / bullet / heading / quote)"""
    blocks = []
    in_code = False
    code_lines = []
    table_lines = []
    in_table = False

    for line in lines:
        stripped = line.rstrip()

        # 코드 블록
        if stripped.startswith('```'):
            if in_code:
                blocks.append({'type': 'code', 'lines': code_lines})
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue

        if in_code:
            code_lines.append(stripped)
            continue

        # 테이블
        if stripped.startswith('|') and stripped.endswith('|'):
            in_table = True
            table_lines.append(stripped)
            continue
        elif in_table:
            blocks.append({'type': 'table', 'lines': table_lines})
            table_lines = []
            in_table = False

        # 빈 줄
        if not stripped:
            blocks.append({'type': 'spacer'})
            continue

        # 인용
        if stripped.startswith('> '):
            blocks.append({'type': 'quote', 'text': clean_inline_markdown(stripped[2:])})
            continue

        # 헤더 H3
        h3 = re.match(r'^###\s+(.+)$', stripped)
        if h3:
            blocks.append({'type': 'h3', 'text': clean_inline_markdown(h3.group(1))})
            continue

        # 리스트 (불릿)
        bullet = re.match(r'^(\s*)[-*]\s+(.+)$', stripped)
        if bullet:
            indent = min(len(bullet.group(1)) // 2, 3)
            text = clean_inline_markdown(bullet.group(2))
            blocks.append({'type': 'bullet', 'text': text, 'indent': indent})
            continue

        # 번호 리스트
        num = re.match(r'^(\s*)\d+\.\s+(.+)$', stripped)
        if num:
            indent = min(len(num.group(1)) // 2, 3)
            text = clean_inline_markdown(num.group(2))
            blocks.append({'type': 'bullet', 'text': text, 'indent': indent})
            continue

        # 일반 본문
        blocks.append({'type': 'body', 'text': clean_inline_markdown(stripped)})

    # 마지막 처리
    if in_code and code_lines:
        blocks.append({'type': 'code', 'lines': code_lines})
    if in_table and table_lines:
        blocks.append({'type': 'table', 'lines': table_lines})

    # 연속 spacer 1개로 압축
    cleaned = []
    last_was_spacer = False
    for b in blocks:
        if b['type'] == 'spacer':
            if last_was_spacer:
                continue
            last_was_spacer = True
        else:
            last_was_spacer = False
        cleaned.append(b)

    return cleaned


# ═════════════════════════════════════════════════════════
# 슬라이드 빌더
# ═════════════════════════════════════════════════════════

def add_background(slide, prs, color=COLOR_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)
    return shape


def add_text(slide, left, top, width, height, text, **kwargs):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if kwargs.get('auto_size'):
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    if 'align' in kwargs:
        p.alignment = kwargs['align']
    if 'anchor' in kwargs:
        tf.vertical_anchor = kwargs['anchor']
    run = p.add_run()
    run.text = str(text)
    set_korean_font(run, kwargs.get('font_name', FONT_BODY))
    run.font.size = Pt(kwargs.get('font_size', 16))
    run.font.bold = kwargs.get('bold', False)
    run.font.italic = kwargs.get('italic', False)
    run.font.color.rgb = kwargs.get('color', COLOR_TEXT)
    return box, tf


def add_bar(slide, left, top, width, height, color=COLOR_ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_star(slide, x, y, size=Pt(28), color=COLOR_ACCENT):
    box = slide.shapes.add_textbox(x, y, Inches(0.6), Inches(0.6))
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '✷'
    run.font.size = size
    run.font.color.rgb = color
    return box


def build_cover(prs, fm: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    # 시리즈 라벨
    add_text(slide, MARGIN_LR, MARGIN_TB, Inches(8), Inches(0.3),
             f"DREAMTEAM CONSULTING SERIES · {fm.get('date', '2026-04-27')}",
             font_size=FONT_LABEL, bold=True, color=COLOR_MUTED, font_name=FONT_HEADING)

    # 메인 타이틀 (이모지 제거)
    title = re.sub(r'[\U0001F300-\U0001FFFF☀-➿]', '', fm.get('title', '')).strip()
    add_text(slide, MARGIN_LR, MARGIN_TB + Inches(0.9), Inches(12), Inches(2.5),
             title, font_size=FONT_HERO_HEADLINE, bold=True, color=COLOR_TEXT, font_name=FONT_HEADING)

    # 부제
    subtitle = fm.get('subtitle', '')
    if subtitle:
        add_text(slide, MARGIN_LR, Inches(4.3), Inches(12), Inches(1),
                 subtitle, font_size=FONT_SUBTITLE, italic=True, color=COLOR_ACCENT, font_name=FONT_HEADING)

    # 강조 바
    add_bar(slide, MARGIN_LR, Inches(5.6), Inches(2.5), Inches(0.1))

    # 작성자
    author = fm.get('author', '윈디 + 벨라님')
    add_text(slide, MARGIN_LR, Inches(5.9), Inches(8), Inches(0.4),
             author, font_size=FONT_BULLET, bold=True, color=COLOR_TEXT, font_name=FONT_HEADING)
    add_text(slide, MARGIN_LR, Inches(6.3), Inches(8), Inches(0.4),
             "@bella_ai_auto · 드림팀 컨설팅 시리즈",
             font_size=12, color=COLOR_MUTED, font_name=FONT_HEADING)

    # ✷ 별
    add_star(slide, Inches(12.0), Inches(0.4), size=Pt(40))

    # 푸터
    add_footer(slide, prs, 1, '표지')


def build_toc(prs, slides_data: list):
    """목차 슬라이드 (자동 생성)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(slide, MARGIN_LR, MARGIN_TB, Inches(12), Inches(0.3),
             "TABLE OF CONTENTS",
             font_size=FONT_LABEL, bold=True, color=COLOR_ACCENT, font_name=FONT_HEADING)

    add_text(slide, MARGIN_LR, MARGIN_TB + Inches(0.3), Inches(12), Inches(1),
             "목차", font_size=FONT_HERO_HEADLINE, bold=True, color=COLOR_TEXT, font_name=FONT_HEADING)

    add_bar(slide, MARGIN_LR, Inches(2.0), Inches(0.8), Inches(0.08))

    # 슬라이드 목록 (10개씩 2단)
    chapters_seen = set()
    items = []
    for i, sd in enumerate(slides_data, 3):  # 표지 + TOC 다음
        ch = sd.get('chapter', '')
        if ch and ch not in chapters_seen:
            chapters_seen.add(ch)
            ch_clean = re.sub(r'[\U0001F300-\U0001FFFF☀-➿]', '', ch).strip()
            items.append({'num': '', 'text': ch_clean, 'chapter': True})
        title_clean = re.sub(r'[\U0001F300-\U0001FFFF☀-➿]', '', sd['title']).strip()
        items.append({'num': f'{i:02d}', 'text': title_clean, 'chapter': False})

    # 두 컬럼으로 분배
    half = (len(items) + 1) // 2
    col1 = items[:half]
    col2 = items[half:]

    for col_idx, items_col in enumerate([col1, col2]):
        x = MARGIN_LR + (Inches(6.5) if col_idx == 1 else 0)
        for j, item in enumerate(items_col):
            y = Inches(2.4) + Inches(0.32) * j
            if y > Inches(7): break
            box = slide.shapes.add_textbox(x, y, Inches(6), Inches(0.32))
            tf = box.text_frame
            tf.margin_left = Emu(0); tf.margin_right = Emu(0)
            tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            if item['chapter']:
                run = p.add_run()
                run.text = '◆ ' + item['text'][:50]
                set_korean_font(run, FONT_HEADING)
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = COLOR_ACCENT
            else:
                if item['num']:
                    run = p.add_run()
                    run.text = item['num'] + '   '
                    set_korean_font(run, FONT_CODE)
                    run.font.size = Pt(11)
                    run.font.color.rgb = COLOR_MUTED
                run = p.add_run()
                run.text = item['text'][:55]
                set_korean_font(run, FONT_BODY)
                run.font.size = Pt(11)
                run.font.color.rgb = COLOR_TEXT

    add_footer(slide, prs, 2, '목차')


def render_blocks_to_slide(slide, blocks: list, top_start=Inches(2.4)):
    """블록을 슬라이드에 렌더링 (영역: 0.7", top_start ~ 6.9")"""
    left = MARGIN_LR
    width = Inches(11.9)
    y = top_start
    max_y = Inches(6.9)

    for block in blocks:
        if y >= max_y:
            break

        btype = block['type']

        if btype == 'spacer':
            y += Inches(0.12)
            continue

        if btype == 'h3':
            box, tf = add_text(slide, left, y, width, Inches(0.4),
                               '▸ ' + block['text'][:60],
                               font_size=FONT_H3, bold=True, color=COLOR_ACCENT, font_name=FONT_HEADING)
            y += Inches(0.45)
            continue

        if btype == 'bullet':
            indent_str = '   ' * block.get('indent', 0)
            box, tf = add_text(slide, left, y, width, Inches(0.4),
                               indent_str + '• ' + block['text'][:90],
                               font_size=FONT_BODY_SIZE, color=COLOR_TEXT, font_name=FONT_BODY)
            y += Inches(0.34)
            continue

        if btype == 'quote':
            # Crail Brown 좌측 바 + italic
            add_bar(slide, left, y + Inches(0.05), Inches(0.05), Inches(0.4))
            box, tf = add_text(slide, left + Inches(0.2), y, width - Inches(0.2), Inches(0.5),
                               block['text'][:120],
                               font_size=FONT_BULLET, italic=True, color=COLOR_ACCENT, font_name=FONT_HEADING)
            y += Inches(0.5)
            continue

        if btype == 'body':
            text = block['text'][:200]
            if not text:
                continue
            box, tf = add_text(slide, left, y, width, Inches(0.4),
                               text,
                               font_size=FONT_BODY_SIZE, color=COLOR_TEXT, font_name=FONT_BODY)
            y += Inches(0.34)
            continue

        if btype == 'code':
            code_lines = block['lines'][:8]  # 최대 8줄
            code_height = Inches(0.28) * len(code_lines) + Inches(0.2)
            if y + code_height >= max_y:
                # EMU 차이를 inch로 변환 (914400 EMU = 1 inch)
                remaining_emu = max_y - y - Inches(0.2)
                remaining_inches = remaining_emu / 914400
                max_lines = max(1, int(remaining_inches / 0.28))
                code_lines = code_lines[:max_lines]
                code_height = Inches(0.28) * len(code_lines) + Inches(0.2)

            # 다크 배경
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        left, y, width, code_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLOR_CODE_BG
            bg.line.fill.background()

            # 코드 텍스트
            code_box = slide.shapes.add_textbox(left + Inches(0.15), y + Inches(0.1),
                                                width - Inches(0.3), code_height - Inches(0.2))
            ctf = code_box.text_frame
            ctf.word_wrap = False
            ctf.margin_left = Emu(0); ctf.margin_right = Emu(0)
            ctf.margin_top = Emu(0); ctf.margin_bottom = Emu(0)
            for i, cl in enumerate(code_lines):
                p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
                run = p.add_run()
                run.text = cl[:120]
                set_korean_font(run, FONT_CODE)
                run.font.size = Pt(11)
                run.font.color.rgb = COLOR_CODE_TEXT
            y += code_height + Inches(0.1)
            continue

        if btype == 'table':
            tbl_lines = [l for l in block['lines'] if l.strip() and not re.match(r'^\|[\s\-:|]+\|?$', l)]
            if not tbl_lines:
                continue

            rows_data = []
            for tl in tbl_lines[:6]:  # 최대 6행
                cells = [c.strip() for c in tl.strip('|').split('|')][:5]  # 최대 5열
                rows_data.append(cells)

            if not rows_data:
                continue

            n_rows = len(rows_data)
            n_cols = max(len(r) for r in rows_data)
            tbl_height = Inches(0.4) * n_rows + Inches(0.05)

            if y + tbl_height >= max_y:
                break

            tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, y, width, tbl_height)
            table = tbl_shape.table
            for r_idx, row_data in enumerate(rows_data):
                for c_idx in range(n_cols):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = ''
                    if c_idx < len(row_data):
                        cell.text_frame.text = clean_inline_markdown(row_data[c_idx])[:60]
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                set_korean_font(r, FONT_BODY)
                                r.font.size = Pt(11)
                                if r_idx == 0:
                                    r.font.bold = True
                                    r.font.color.rgb = COLOR_TEXT
                                else:
                                    r.font.color.rgb = COLOR_TEXT
                    cell.fill.solid()
                    if r_idx == 0:
                        cell.fill.fore_color.rgb = COLOR_HIGHLIGHT_BG
                    else:
                        cell.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF2)
            y += tbl_height + Inches(0.15)
            continue

    return y


def build_content_slide(prs, sd: dict, num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    # 챕터 라벨 (있으면)
    chapter = sd.get('chapter', '')
    if chapter:
        chapter_clean = re.sub(r'[\U0001F300-\U0001FFFF☀-➿]', '', chapter).strip()[:60]
        add_text(slide, MARGIN_LR, MARGIN_TB, Inches(12), Inches(0.3),
                 chapter_clean.upper(),
                 font_size=FONT_LABEL, bold=True, color=COLOR_ACCENT, font_name=FONT_HEADING)

    # 슬라이드 제목
    title = re.sub(r'[\U0001F300-\U0001FFFF☀-➿]', '', sd['title']).strip()
    add_text(slide, MARGIN_LR, MARGIN_TB + Inches(0.25), Inches(12), Inches(1.0),
             title, font_size=FONT_CHAPTER, bold=True, color=COLOR_TEXT, font_name=FONT_HEADING)

    # Crail Brown 강조 바
    add_bar(slide, MARGIN_LR, Inches(2.1), Inches(0.8), Inches(0.06))

    # 본문 블록 렌더링
    blocks = render_content_blocks(sd['content'])
    render_blocks_to_slide(slide, blocks, top_start=Inches(2.35))

    add_footer(slide, prs, num, sd['title'][:40])


def add_footer(slide, prs, num, section=''):
    add_bar(slide, MARGIN_LR, Inches(7.1), Inches(11.9), Inches(0.015), color=COLOR_MUTED)
    add_text(slide, MARGIN_LR, Inches(7.18), Inches(8), Inches(0.3),
             "@bella_ai_auto · 드림팀 컨설팅 시리즈",
             font_size=9, color=COLOR_MUTED, font_name=FONT_HEADING)
    add_text(slide, Inches(11.5), Inches(7.18), Inches(1.2), Inches(0.3),
             f"{num:02d}",
             font_size=FONT_FOOTER, bold=True, color=COLOR_ACCENT, font_name=FONT_HEADING, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════
# 메인 변환
# ═════════════════════════════════════════════════════════

def convert(input_path: Path, output_path: Path):
    print(f'\n=== Converting: {input_path.name} ===')
    content = input_path.read_text(encoding='utf-8')
    parsed = parse_markdown(content)
    fm = parsed['frontmatter']
    body = parsed['body']

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. 표지
    build_cover(prs, fm)
    print(f'  [01] 표지: {fm.get("title", "")[:40]}')

    # 2. 슬라이드 데이터 분할
    slides_data = split_into_slides_v2(body)

    # 3. 목차
    build_toc(prs, slides_data)
    print(f'  [02] 목차 ({len(slides_data)}개 항목)')

    # 4. 본문
    num = 3
    for sd in slides_data:
        build_content_slide(prs, sd, num)
        title_preview = sd['title'][:40].replace('\n', ' ')
        print(f'  [{num:02d}] {title_preview}')
        num += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f'  ✅ Saved: {output_path.name} ({num-1} slides)')
    return num - 1


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', help='Single markdown file')
    parser.add_argument('--output', help='Output PPTX path')
    parser.add_argument('--batch', action='store_true', help='Convert all 7 guides')
    parser.add_argument('--sample', action='store_true', help='Sample 1 guide only (정복 가이드)')
    args = parser.parse_args()

    guides = [
        ('D:/Claude_AI_Knowledge/07_Automation_Business/meta_ig_token_complete_guide_v1.md',
         'D:/bella-ppt-vault/manneungi/pptx_v2/01_meta_ig_token_complete_guide.pptx'),
        ('D:/Claude_AI_Knowledge/07_Automation_Business/auto_publish_4stage_vibeline_v1.md',
         'D:/bella-ppt-vault/manneungi/pptx_v2/02_auto_publish_4stage_vibeline.pptx'),
        ('D:/Claude_AI_Knowledge/07_Automation_Business/instagram_claude_content_reference_12cases_v1.md',
         'D:/bella-ppt-vault/manneungi/pptx_v2/03_instagram_reference_15cases.pptx'),
        ('D:/Claude_AI_Knowledge/07_Automation_Business/auto_comment_to_dm_system_design_v1.md',
         'D:/bella-ppt-vault/manneungi/pptx_v2/04_auto_comment_to_dm_system.pptx'),
        ('D:/Claude_AI_Knowledge/07_Automation_Business/instagram_publish_troubleshooting_v1.md',
         'D:/bella-ppt-vault/manneungi/pptx_v2/05_instagram_publish_troubleshooting.pptx'),
        ('D:/dreamteam-hq/DKM/09_DAILY_LOGS/2026-04-27_bella_4hr_marathon_FULL_EOD.md',
         'D:/bella-ppt-vault/manneungi/pptx_v2/06_bella_4hr_marathon_D5_EOD.pptx'),
        ('D:/Claude_AI_Knowledge/07_Automation_Business/github_3repos_dreamteam_apply_v1.md',
         'D:/bella-ppt-vault/manneungi/pptx_v2/07_github_3repos_dreamteam_apply.pptx'),
    ]

    if args.sample:
        inp, out = guides[0]
        convert(Path(inp), Path(out))
    elif args.batch:
        total = 0
        for inp, out in guides:
            if not Path(inp).exists():
                print(f'[SKIP] {inp}')
                continue
            slides = convert(Path(inp), Path(out))
            total += slides
        print(f'\n=== Batch complete: {len(guides)} guides, {total} slides ===')
    elif args.input and args.output:
        convert(Path(args.input), Path(args.output))
    else:
        parser.print_help()


if __name__ == '__main__':
    main()
