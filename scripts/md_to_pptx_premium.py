"""
🎨 마크다운 → PPTX 전문가급 자동 변환 (Anthropic 브랜드 톤)

특징:
- Anthropic 공식 컬러 (Crail Brown #CC785C + Manilla Cream #F0EEE6 + Slate #191919)
- 한글 폰트 (Pretendard 우선, 폴백 맑은 고딕)
- 코드 블록 + 테이블 + 인용구 보존
- 표지 슬라이드 (제목 + 부제 + 날짜 + 저자)
- 챕터 구분 슬라이드 (Crail Brown 강조)
- 본문 가독성 (제목 32pt + 본문 16pt + 코드 13pt monospace)
- 푸터 + 슬라이드 번호

사용법:
  python md_to_pptx_premium.py --input <md_file> --output <pptx_file>
  python md_to_pptx_premium.py --batch       # 6개 가이드 일괄 변환
"""
import argparse
import re
import sys
import io
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# ═════════════════════════════════════════════════════════
# 🎨 Anthropic 공식 브랜드 (https://claude.com/brand)
# ═════════════════════════════════════════════════════════
COLOR_BG = RGBColor(0xF0, 0xEE, 0xE6)        # Manilla Cream
COLOR_TEXT = RGBColor(0x19, 0x19, 0x19)      # Slate near-black
COLOR_ACCENT = RGBColor(0xCC, 0x78, 0x5C)    # Crail Brown
COLOR_MUTED = RGBColor(0x80, 0x70, 0x60)     # warm gray
COLOR_CODE_BG = RGBColor(0x1A, 0x1A, 0x1A)   # dark code bg
COLOR_CODE_TEXT = RGBColor(0xF0, 0xEE, 0xE6) # cream on dark
COLOR_CODE_ACCENT = RGBColor(0xCC, 0x78, 0x5C)
COLOR_HIGHLIGHT = RGBColor(0xF5, 0xE0, 0xD0) # subtle highlight

FONT_HEADING = '맑은 고딕'  # 한글 + 영문 안전 폴백
FONT_BODY = '맑은 고딕'
FONT_CODE = 'Consolas'

# 슬라이드 크기 (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ═════════════════════════════════════════════════════════
# 마크다운 파서
# ═════════════════════════════════════════════════════════

def parse_markdown(content: str) -> dict:
    """프론트매터 + 본문 파싱"""
    fm = {}
    body = content

    # YAML 프론트매터
    fm_match = re.match(r'^---\n(.*?)\n---\n(.*)', content, re.DOTALL)
    if fm_match:
        fm_text = fm_match.group(1)
        body = fm_match.group(2)
        for line in fm_text.splitlines():
            if ':' in line:
                k, v = line.split(':', 1)
                fm[k.strip()] = v.strip()

    return {'frontmatter': fm, 'body': body}


def split_into_slides(body: str) -> list:
    """마크다운 본문을 슬라이드 단위로 분할
    - # H1 → 표지 또는 챕터 구분
    - ## H2 → 새 슬라이드 (메인)
    - ### H3 → 새 슬라이드 (서브)
    """
    slides = []
    lines = body.split('\n')
    current = {'title': None, 'level': None, 'content': []}

    for line in lines:
        # H1 (챕터 구분)
        h1 = re.match(r'^# (.+)$', line)
        # H2 (메인)
        h2 = re.match(r'^## (.+)$', line)
        # H3 (서브)
        h3 = re.match(r'^### (.+)$', line)

        if h1:
            if current['title']:
                slides.append(current)
            current = {'title': h1.group(1).strip(), 'level': 1, 'content': []}
        elif h2:
            if current['title']:
                slides.append(current)
            current = {'title': h2.group(1).strip(), 'level': 2, 'content': []}
        elif h3:
            if current['title']:
                slides.append(current)
            current = {'title': h3.group(1).strip(), 'level': 3, 'content': []}
        else:
            current['content'].append(line)

    if current['title']:
        slides.append(current)

    return slides


# ═════════════════════════════════════════════════════════
# 슬라이드 빌더
# ═════════════════════════════════════════════════════════

def add_background(slide, prs, color=COLOR_BG):
    """배경색 채우기"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # 테두리 제거
    # 맨 뒤로 보내기
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)
    return shape


def add_text_box(slide, left, top, width, height, text, **kwargs):
    """텍스트 박스 추가 (옵션: font_size, bold, color, font_name, align)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    if 'align' in kwargs:
        p.alignment = kwargs['align']

    run = p.add_run()
    run.text = str(text)
    run.font.name = kwargs.get('font_name', FONT_BODY)
    run.font.size = Pt(kwargs.get('font_size', 16))
    run.font.bold = kwargs.get('bold', False)
    run.font.color.rgb = kwargs.get('color', COLOR_TEXT)
    return box


def add_accent_bar(slide, left, top, width, height, color=COLOR_ACCENT):
    """Crail Brown 강조 바"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_star_motif(slide, x, y, size=Inches(0.5), color=COLOR_ACCENT):
    """Anthropic ✷ 4-pointed star motif"""
    # Use a simple text-based asterisk in Crail Brown
    box = slide.shapes.add_textbox(x, y, size, size)
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '✷'
    run.font.size = Pt(28)
    run.font.color.rgb = color
    return box


def build_cover_slide(prs, fm: dict):
    """표지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, prs)

    # 좌측 상단 라벨
    add_text_box(
        slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.3),
        f"DREAMTEAM CONSULTING SERIES · {fm.get('date', '')}",
        font_size=10, bold=True, color=COLOR_MUTED, font_name=FONT_HEADING
    )

    # 타이틀
    title = fm.get('title', 'Untitled').replace('🎓', '').replace('🚀', '').replace('📚', '').replace('🤖', '').replace('🚨', '').replace('🏆', '').strip()
    add_text_box(
        slide, Inches(0.7), Inches(1.5), Inches(12), Inches(2.5),
        title,
        font_size=44, bold=True, color=COLOR_TEXT, font_name=FONT_HEADING
    )

    # 부제
    subtitle = fm.get('subtitle', '')
    if subtitle:
        add_text_box(
            slide, Inches(0.7), Inches(4.3), Inches(12), Inches(1),
            subtitle,
            font_size=18, color=COLOR_ACCENT, font_name=FONT_HEADING
        )

    # Crail Brown 강조 바
    add_accent_bar(slide, Inches(0.7), Inches(5.6), Inches(2), Inches(0.08))

    # 작성자 + 날짜
    author = fm.get('author', '윈디 + 벨라님')
    add_text_box(
        slide, Inches(0.7), Inches(5.9), Inches(8), Inches(0.4),
        f"{author}",
        font_size=14, bold=True, color=COLOR_TEXT, font_name=FONT_HEADING
    )
    add_text_box(
        slide, Inches(0.7), Inches(6.3), Inches(8), Inches(0.4),
        f"@bella_ai_auto · 드림팀 컨설팅 시리즈",
        font_size=12, color=COLOR_MUTED, font_name=FONT_HEADING
    )

    # ✷ 모티프
    add_star_motif(slide, Inches(12.0), Inches(0.5))

    # 푸터
    add_footer(slide, prs, 1, '표지')


def build_chapter_slide(prs, title, num):
    """챕터 구분 슬라이드 (H1)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    # 챕터 번호
    add_text_box(
        slide, Inches(0.7), Inches(2.5), Inches(12), Inches(0.5),
        f"CHAPTER",
        font_size=14, bold=True, color=COLOR_ACCENT, font_name=FONT_HEADING
    )

    # 큰 챕터 제목
    clean_title = clean_emoji_for_pptx(title)
    add_text_box(
        slide, Inches(0.7), Inches(3.0), Inches(12), Inches(2),
        clean_title,
        font_size=54, bold=True, color=COLOR_TEXT, font_name=FONT_HEADING
    )

    # Crail Brown 가로 라인
    add_accent_bar(slide, Inches(0.7), Inches(5.0), Inches(3), Inches(0.1))

    # ✷ 모티프
    add_star_motif(slide, Inches(11.5), Inches(3.2), size=Inches(0.8))

    add_footer(slide, prs, num, 'CHAPTER')


def clean_emoji_for_pptx(text: str) -> str:
    """PPTX에서 일부 깨지는 이모지 정리"""
    return text


def build_content_slide(prs, slide_data: dict, num):
    """일반 본문 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    # 상단 작은 라벨 (레벨에 따라)
    level_label = {1: 'CHAPTER', 2: 'TOPIC', 3: 'DETAIL'}.get(slide_data['level'], 'CONTENT')
    add_text_box(
        slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.3),
        level_label,
        font_size=10, bold=True, color=COLOR_ACCENT, font_name=FONT_HEADING
    )

    # 슬라이드 제목
    title = clean_emoji_for_pptx(slide_data['title'])
    title_size = 32 if slide_data['level'] == 2 else 28
    add_text_box(
        slide, Inches(0.7), Inches(0.8), Inches(12), Inches(1),
        title,
        font_size=title_size, bold=True, color=COLOR_TEXT, font_name=FONT_HEADING
    )

    # Crail Brown 강조 바 (제목 아래)
    add_accent_bar(slide, Inches(0.7), Inches(2.0), Inches(0.8), Inches(0.06))

    # 본문 (마크다운 → 슬라이드 텍스트)
    content_lines = slide_data['content']
    rendered = render_content_blocks(content_lines)

    # 본문 영역
    body_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.3), Inches(12), Inches(4.7)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    first = True
    for block in rendered:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        text = block['text']
        run = p.add_run()
        run.text = text

        if block['type'] == 'code':
            run.font.name = FONT_CODE
            run.font.size = Pt(13)
            run.font.color.rgb = COLOR_CODE_TEXT
            # 배경은 직접 표현 어려워서 색상만
        elif block['type'] == 'heading':
            run.font.name = FONT_HEADING
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = COLOR_ACCENT
        elif block['type'] == 'bullet':
            run.font.name = FONT_BODY
            run.font.size = Pt(15)
            run.font.color.rgb = COLOR_TEXT
            p.level = block.get('indent', 0)
        elif block['type'] == 'quote':
            run.font.name = FONT_BODY
            run.font.size = Pt(15)
            run.font.italic = True
            run.font.color.rgb = COLOR_ACCENT
        elif block['type'] == 'table':
            run.font.name = FONT_CODE
            run.font.size = Pt(11)
            run.font.color.rgb = COLOR_TEXT
        else:  # body
            run.font.name = FONT_BODY
            run.font.size = Pt(15)
            run.font.color.rgb = COLOR_TEXT

    add_footer(slide, prs, num, slide_data['title'][:40])


def render_content_blocks(lines: list) -> list:
    """마크다운 라인들을 블록으로 변환"""
    blocks = []
    in_code = False
    code_lines = []

    for line in lines:
        stripped = line.rstrip()

        # 코드 블록
        if stripped.startswith('```'):
            if in_code:
                blocks.append({'type': 'code', 'text': '\n'.join(code_lines)})
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue

        if in_code:
            code_lines.append(stripped)
            continue

        # 빈 줄
        if not stripped:
            blocks.append({'type': 'body', 'text': ''})
            continue

        # 인용
        if stripped.startswith('> '):
            blocks.append({'type': 'quote', 'text': stripped[2:]})
            continue

        # 테이블 (간단 처리)
        if stripped.startswith('|') and '|' in stripped[1:]:
            blocks.append({'type': 'table', 'text': stripped})
            continue

        # 리스트
        bullet_match = re.match(r'^(\s*)[-*]\s+(.+)$', stripped)
        if bullet_match:
            indent_chars = len(bullet_match.group(1))
            indent = min(indent_chars // 2, 4)
            text = bullet_match.group(2)
            blocks.append({'type': 'bullet', 'text': f'• {text}', 'indent': indent})
            continue

        # 번호 리스트
        num_match = re.match(r'^(\s*)\d+\.\s+(.+)$', stripped)
        if num_match:
            indent_chars = len(num_match.group(1))
            indent = min(indent_chars // 2, 4)
            text = num_match.group(2)
            blocks.append({'type': 'bullet', 'text': f'• {text}', 'indent': indent})
            continue

        # 굵은 강조 라인 (전체가 ** **)
        if stripped.startswith('**') and stripped.endswith('**'):
            blocks.append({'type': 'heading', 'text': stripped[2:-2]})
            continue

        # 일반 본문 (인라인 마크다운 정리)
        # **bold** → bold (강조 효과는 PPTX에서 한 run이라 단순화)
        cleaned = re.sub(r'\*\*(.+?)\*\*', r'\1', stripped)
        cleaned = re.sub(r'\*(.+?)\*', r'\1', cleaned)
        cleaned = re.sub(r'`([^`]+)`', r'\1', cleaned)

        blocks.append({'type': 'body', 'text': cleaned})

    if in_code and code_lines:
        blocks.append({'type': 'code', 'text': '\n'.join(code_lines)})

    return blocks


def add_footer(slide, prs, num, section_title=''):
    """푸터: 좌측 시리즈명, 우측 슬라이드 번호"""
    # 하단 가는 라인
    add_accent_bar(slide, Inches(0.7), Inches(7.05), Inches(11.9), Inches(0.02), color=COLOR_MUTED)

    # 좌측: @bella_ai_auto · 드림팀 컨설팅
    add_text_box(
        slide, Inches(0.7), Inches(7.15), Inches(8), Inches(0.3),
        f"@bella_ai_auto · 드림팀 컨설팅 시리즈",
        font_size=9, color=COLOR_MUTED, font_name=FONT_HEADING
    )

    # 우측: 슬라이드 번호
    add_text_box(
        slide, Inches(11.5), Inches(7.15), Inches(1.2), Inches(0.3),
        f"{num:02d}",
        font_size=10, bold=True, color=COLOR_ACCENT, font_name=FONT_HEADING, align=PP_ALIGN.RIGHT
    )


# ═════════════════════════════════════════════════════════
# 메인 변환 로직
# ═════════════════════════════════════════════════════════

def convert(input_path: Path, output_path: Path):
    print(f'\n=== Converting: {input_path.name} ===')
    content = input_path.read_text(encoding='utf-8')
    parsed = parse_markdown(content)
    fm = parsed['frontmatter']
    body = parsed['body']

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. 표지
    build_cover_slide(prs, fm)
    print(f'  [01] 표지: {fm.get("title", "Untitled")[:40]}')

    # 2. 본문 슬라이드들
    slides_data = split_into_slides(body)
    num = 2
    for sd in slides_data:
        if sd['level'] == 1:
            # H1은 챕터 구분
            build_chapter_slide(prs, sd['title'], num)
            print(f'  [{num:02d}] CHAPTER: {sd["title"][:40]}')
            num += 1
            # 본문 내용도 추가 슬라이드
            if any(l.strip() for l in sd['content']):
                build_content_slide(prs, sd, num)
                print(f'  [{num:02d}] {sd["title"][:40]}')
                num += 1
        else:
            build_content_slide(prs, sd, num)
            print(f'  [{num:02d}] {sd["title"][:40]}')
            num += 1

    # 저장
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f'  ✅ Saved: {output_path} ({num-1} slides)')
    return num - 1


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', help='Single markdown file')
    parser.add_argument('--output', help='Output PPTX path')
    parser.add_argument('--batch', action='store_true', help='Convert all 6 guides')
    args = parser.parse_args()

    if args.batch:
        # 6개 가이드 일괄 변환
        guides = [
            ('D:/Claude_AI_Knowledge/07_Automation_Business/meta_ig_token_complete_guide_v1.md',
             'D:/bella-ppt-vault/manneungi/pptx/01_meta_ig_token_complete_guide_v1.pptx'),
            ('D:/Claude_AI_Knowledge/07_Automation_Business/auto_publish_4stage_vibeline_v1.md',
             'D:/bella-ppt-vault/manneungi/pptx/02_auto_publish_4stage_vibeline_v1.pptx'),
            ('D:/Claude_AI_Knowledge/07_Automation_Business/instagram_claude_content_reference_12cases_v1.md',
             'D:/bella-ppt-vault/manneungi/pptx/03_instagram_reference_15cases_v1.pptx'),
            ('D:/Claude_AI_Knowledge/07_Automation_Business/auto_comment_to_dm_system_design_v1.md',
             'D:/bella-ppt-vault/manneungi/pptx/04_auto_comment_to_dm_system_v1.pptx'),
            ('D:/Claude_AI_Knowledge/07_Automation_Business/instagram_publish_troubleshooting_v1.md',
             'D:/bella-ppt-vault/manneungi/pptx/05_instagram_publish_troubleshooting_v1.pptx'),
            ('D:/dreamteam-hq/DKM/09_DAILY_LOGS/2026-04-27_bella_4hr_marathon_FULL_EOD.md',
             'D:/bella-ppt-vault/manneungi/pptx/06_bella_4hr_marathon_D5_EOD.pptx'),
        ]
        total = 0
        for inp, out in guides:
            inp_path = Path(inp)
            out_path = Path(out)
            if not inp_path.exists():
                print(f'[SKIP] Not found: {inp}')
                continue
            slides = convert(inp_path, out_path)
            total += slides
        print(f'\n=== Batch complete: {len(guides)} guides, {total} total slides ===')
    elif args.input and args.output:
        convert(Path(args.input), Path(args.output))
    else:
        parser.print_help()


if __name__ == '__main__':
    main()
